#!/usr/bin/python
#encoding=utf-8
from googletrans import Translator
import os
import tkinter.filedialog
import tkinter.messagebox
from pptx import Presentation
import time
import re

gtrans = Translator(
    service_urls=['translate.google.cn'],
    proxies={
#        'https': '183.148.150.160:9999',
    })

gtrans.translate("Apple", src='en', dest='zh-cn')

def do_translate(instr):
    while(True):
        try:
            result = gtrans.translate(instr, src='en', dest='zh-cn')
            return result.text
        except:
            print("Google服务器连接失败，即将重试")
            time.sleep(5)

def translate_list(strlist):
    outlist = []
    trstr = strlist[0]
    gp_num = 1
    for i in range(len(strlist)):
        if(i > 0):
            trstr = trstr + "\n\n\n\n\n\n" + strlist[i]
            gp_num = gp_num + 1
        if((len(trstr) > 1000) or (i == len(strlist)-1)):
            trans_str = do_translate(trstr)
            split_trans_str = re.split('\n\n\n\n\n\n', trans_str)
            outlist.extend(split_trans_str)
            trstr = ''
    return(outlist)


def do_parse_and_translate_pptx(start_slide, end_slide, infile, outfile, trans_notes=True):
    try:
        prs = Presentation(infile)
    except IOError:
        tkinter.messagebox.showerror('错误', "无法读取文件%s"%infile)
        exit(1)
    num_slides = len(prs.slides)
    if(end_slide == -1): end_slide = num_slides - 1
    fail_page = []
    for i in range(num_slides):
        if(i < start_slide): continue
        if(i > end_slide): break
        slide = prs.slides[i]
        print("正在翻译第%0d/%0d页"%(i,num_slides))
        text_list = []
        for shape in slide.shapes:
            if(shape.has_text_frame):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if(re.match(r'^\s*$', run.text, flags=re.M)):
                            text_list.append('12345678')
                        else:
                            text_list.append(run.text)
            if(shape.has_table):
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if(re.match(r'^\s*$', run.text, flags=re.M)):
                                text_list.append('12345678')
                            else:
                                text_list.append(run.text)
        if(trans_notes and slide.has_notes_slide):
            notes_slide = slide.notes_slide
            for paragraph in notes_slide.notes_text_frame.paragraphs:
                for run in paragraph.runs:
                    if(re.match(r'^\s*$', run.text, flags=re.M)):
                        text_list.append('12345678')
                    else:
                        text_list.append(run.text)
        if(len(text_list) > 0):
            trans_text_list = translate_list(text_list)
            if(len(trans_text_list) != len(text_list)):
                print("在翻译第%0d页中发生错误"%i)
                print("英文列表包含%0d个元素"%len(text_list))
                print(text_list)
                print("翻译后列表包含%0d个元素"%len(trans_text_list))
                print(trans_text_list)
                print("将跳过该页后继续")
                fail_page.append(str(i+1))
                continue
            for shape in slide.shapes:
                if(shape.has_text_frame):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            tmptext = trans_text_list.pop(0)
                            if(tmptext == '12345678'):
                                run.text = ''
                            else:
                                run.text = tmptext
                if(shape.has_table):
                    for cell in shape.table.iter_cells():
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                tmptext = trans_text_list.pop(0)
                                if(tmptext == '12345678'):
                                    run.text = ''
                                else:
                                    run.text = tmptext
            if(trans_notes and slide.has_notes_slide):
                notes_slide = slide.notes_slide
                for paragraph in notes_slide.notes_text_frame.paragraphs:
                    for run in paragraph.runs:
                        tmptext = trans_text_list.pop(0)
                        if(tmptext == '12345678'):
                            run.text = ''
                        else:
                            run.text = tmptext
        while(True):
            try:
                prs.save(outfile)
                print("翻译结果写入%s"%outfile)
                if(i > 0 and i % 10 == 0):
                    wait_time = 5
                    print("等待%0d秒后继续"%wait_time)
                    time.sleep(wait_time)
                break
            except IOError:
                retry = tkinter.messagebox.askretrycancel("错误", "写入%s失败，请关闭文档后重试"%outfile)
                if(retry==False):
                    break
    total_trans_num = end_slide - start_slide + 1
    tkinter.messagebox.showinfo("提示","翻译结束，结果生成在%s\n共翻译%0d页，成功%0d页，失败%0d页\n错误页码:%s"%(
        outfile, total_trans_num, total_trans_num-len(fail_page), len(fail_page), ' '.join(fail_page)))

default_dir = os.getcwd()
infile = tkinter.filedialog.askopenfilename(title=u"选择文件", initialdir=(os.path.expanduser(default_dir)),
filetypes=[("PowerPoint","*.pptx")])
outfile = ''
if(infile != ''):
    outfile = infile.replace(".pptx","_zh.pptx")
    yesno = tkinter.messagebox.askyesno("问题","是否使用\'%s\'作为输出文件名"%outfile)
    if(yesno == False):
        outfile = tkinter.filedialog.asksaveasfilename(title=u"请输入文件名", initialdir=(os.path.expanduser(default_dir)),
filetypes=[("PowerPointFile","*.pptx")], defaultextension=".pptx")
else:
    tkinter.messagebox.showerror("错误","未选择任何文件!")
    exit()
if(outfile == ''):
    tkinter.messagebox.showerror("错误", "未指定输出文件")
    exit()

yesno = tkinter.messagebox.askyesno("问题","是否指定翻译页数范围")
start = 0
end = -1
if(yesno == True):
    print("从第几页开始翻译:")
    while(True):
        try:
            start = int(input())
            if(start > 0):
                start = start - 1
            break
        except:
            print("输入页码非法")
    print("到第几页结束(输入-1表示末页):")
    while(True):
        try:
            end = int(input())
            if(end != -1):
                end = end - 1
            break
        except:
            print("输入页码非法")
need_trans_notes = tkinter.messagebox.askyesno("问题","是否翻译备注内容")
do_parse_and_translate_pptx(start, end, infile, outfile, need_trans_notes)

